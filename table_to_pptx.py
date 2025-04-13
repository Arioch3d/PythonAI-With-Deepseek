#Convert a markdown table to excel using pptx
from pptx import Presentation

markdown_table = """| Name | Age | Occupation | Favorite Hobby |
|--------------|--------|----------------|----------------|
| Alice        | 28     | Engineer       | Painting       |   
| Bob          | 35     | Doctor         | Cycling        |
| Charlie      | 42     | Teacher        | Hiking         |
| David        | 30     | Photographer   | Chess          |"""

def create_ppt(markdown_table, ppt_file_name):
    rows = markdown_table.split("\n")
    spreadsheet=[]
    for index,row in enumerate(rows):
        if index==1:
            continue
        split_rows = row.split("|")
        spreadsheet.append(split_rows[1:-1])

    headers = spreadsheet[0]
    dict_list = [dict(zip(headers,row))for row in spreadsheet[1:]]
    presentation = Presentation()

    for row in dict_list:
        slide_layout = presentation.slide_layouts[1]
        
        slide = presentation.slides.add_slide(slide_layout)
        first_col = headers[0] #Grabs the name of the first column
        slide.placeholders[0].text = row[first_col] #Value of the first column header
        content_text = "\n".join(f"{key.strip('*').strip}:{value}" for key,value in row.items()) #joins the column header and value on 1 line.
        slide.placeholders[1].text = content_text

    presentation.save(f"{ppt_file_name}.pptx")

create_ppt(markdown_table, "table-to-powerpoint")