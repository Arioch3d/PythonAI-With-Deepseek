#1. import library
from pptx import Presentation

#2. create presentation
presentation = Presentation()
#3. Create slide
#3.1 layut selection
slide_layout = presentation.slide_layouts[1]
#3.2 slide creation
slide = presentation.slides.add_slide(slide_layout)
#4. add title
slide.placeholders[0].text = "My Title"
#5. add content
slide.placeholders[1].text = "Apple\nOranges"
#6. save presentation
presentation.save("sample.pptx")